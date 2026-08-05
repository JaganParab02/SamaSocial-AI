"""
PPTParser using python-pptx.
Extracts slide text, slide numbers, speaker notes placeholders, title, and slide metadata.
"""

import os
from typing import Any, Dict, List, Optional
from app.parsers.base import BaseParser
from app.models.document import ParserResult, SourceType
from app.core.logger import get_logger

logger = get_logger(__name__)

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    logger.error("python-pptx package is not installed.")


class PPTParser(BaseParser):
    """
    Parser for PowerPoint presentation slides (.ppt, .pptx).
    Extracts structured content per slide to assist learning Q&A and course structure sequencing.
    """

    def extract(self, source: str, **kwargs: Any) -> ParserResult:
        """
        Extract slide text, sequence numbers, speaker notes, and presentation titles.
        
        Args:
            source: Path to PowerPoint presentation file on disk.
        """
        if not HAS_PPTX:
            raise RuntimeError("python-pptx library required for presentation parsing is uninstalled.")
        if not os.path.exists(source):
            raise FileNotFoundError(f"Presentation document file not found at: {source}")

        logger.info("Extracting presentation slides from: %s", source)
        try:
            prs = Presentation(source)
            extracted_slides: List[Dict[str, Any]] = []
            full_text = []
            doc_title = os.path.basename(source)

            for slide_idx, slide in enumerate(prs.slides, start=1):
                slide_texts = []
                # Extract text from all shapes in slide layout
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_val = shape.text.strip()
                        if text_val:
                            slide_texts.append(text_val)

                # Capture presentation title from first slide title box if available
                if slide_idx == 1 and slide.shapes.title and slide.shapes.title.text:
                    doc_title = slide.shapes.title.text.strip()

                # Extract speaker notes placeholder or text if attached to slide
                notes_text = ""
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                notes_placeholder = notes_text if notes_text else "Speaker notes placeholder - No notes attached to slide"

                combined_slide_text = "\n".join(slide_texts)
                if combined_slide_text:
                    formatted_text = f"--- Slide {slide_idx} ---\n{combined_slide_text}"
                    if notes_text:
                        formatted_text += f"\n[Speaker Notes: {notes_text}]"
                    
                    full_text.append(formatted_text)
                    extracted_slides.append({
                        "slide_number": slide_idx,
                        "text": combined_slide_text,
                        "notes_placeholder": notes_placeholder,
                    })

            meta = self.metadata(prs, title=doc_title, count=len(prs.slides))
            combined_text = "\n\n".join(full_text)

            return ParserResult(
                source_name=doc_title or os.path.basename(source),
                source_type=SourceType.PPT,
                text_content=combined_text,
                metadata=meta,
                structured_items=extracted_slides,
                summary_placeholder=self.summary(combined_text, slides=len(extracted_slides))
            )
        except Exception as e:
            logger.error("Failed extracting presentation structure from '%s': %s", source, str(e))
            raise RuntimeError(f"PowerPoint extraction error: {str(e)}")

    def metadata(self, source: Any, **kwargs: Any) -> Dict[str, Any]:
        """Extract core presentation parameters and slide counts."""
        try:
            prs = source if isinstance(source, Presentation) else Presentation(source)
            meta = {
                "title": kwargs.get("title", ""),
                "slide_count": len(prs.slides),
                "slide_width_in_inches": round(prs.slide_width / 914400, 2) if prs.slide_width else None,
                "slide_height_in_inches": round(prs.slide_height / 914400, 2) if prs.slide_height else None,
            }
            return {k: v for k, v in meta.items() if v is not None}
        except Exception as e:
            logger.warning("Could not parse PPT metadata: %s", str(e))
            return {"slide_count": kwargs.get("count", "unknown")}

    def summary(self, text_content: str, **kwargs: Any) -> Optional[str]:
        """Generate high-level presentation digest string."""
        slides = kwargs.get("slides", 1)
        word_count = len(text_content.split())
        return f"PowerPoint Presentation containing {slides} slides with ~{word_count} words total."

    @classmethod
    def supported_extensions(cls) -> List[str]:
        return ["ppt", "pptx"]
